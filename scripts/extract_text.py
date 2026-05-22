"""Extract text from PDF/docx/pptx in cours-raw → cours-text/<same-relpath>.md."""
from __future__ import annotations

import sys
from pathlib import Path

import pdfplumber
from docx import Document
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
RAW = ROOT / "cours-raw"
TEXT = ROOT / "cours-text"


def extract_pdf(p: Path) -> str:
    out = []
    with pdfplumber.open(p) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            t = page.extract_text() or ""
            out.append(f"## Page {i}\n\n{t.strip()}\n")
    return "\n".join(out)


def extract_docx(p: Path) -> str:
    doc = Document(p)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def extract_pptx(p: Path) -> str:
    # .ppsx (PowerPoint Show) has a different content type and python-pptx refuses to open it.
    # Workaround: rewrite the inner content type to the presentation one in a temp copy.
    src = p
    if p.suffix.lower() == ".ppsx":
        import shutil, tempfile, zipfile, re
        tmp = Path(tempfile.mkstemp(suffix=".pptx")[1])
        shutil.copy(p, tmp)
        # rewrite [Content_Types].xml inside the zip
        with zipfile.ZipFile(tmp, "r") as zin:
            data = {n: zin.read(n) for n in zin.namelist()}
        ct = data["[Content_Types].xml"].decode("utf-8")
        ct = ct.replace("presentationml.slideshow.main+xml", "presentationml.presentation.main+xml")
        data["[Content_Types].xml"] = ct.encode("utf-8")
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for n, d in data.items():
                zout.writestr(n, d)
        src = tmp
    prs = Presentation(src)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        chunks = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        chunks.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    chunks.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunks.append(f"\n**Notes:** {notes}")
        out.append("\n\n".join(chunks))
    return "\n\n---\n\n".join(out)


def extract_txt(p: Path) -> str:
    return p.read_text(errors="replace")


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".ppsx": extract_pptx,
    ".txt": extract_txt,
}


def process(rel: str) -> bool:
    src = RAW / rel
    if not src.exists():
        print(f"[extract] missing source {rel}", file=sys.stderr)
        return False
    fn = EXTRACTORS.get(src.suffix.lower())
    if not fn:
        print(f"[extract] unsupported ext {src.suffix} for {rel}", file=sys.stderr)
        return False
    try:
        body = fn(src)
    except Exception as e:
        print(f"[extract] failed {rel}: {e}", file=sys.stderr)
        return False
    dst = TEXT / Path(rel).with_suffix(".md")
    dst.parent.mkdir(parents=True, exist_ok=True)
    header = f"<!-- source: {rel} -->\n# {Path(rel).stem}\n\n"
    dst.write_text(header + body)
    print(f"[extract] {rel} → {dst.relative_to(ROOT)}")
    return True


def main() -> int:
    todo = ROOT / "scripts" / ".todo_extract.txt"
    if len(sys.argv) > 1 and sys.argv[1] == "--all":
        targets = [
            str(p.relative_to(RAW))
            for p in RAW.rglob("*")
            if p.is_file() and p.suffix.lower() in EXTRACTORS and "_staging" not in p.parts
        ]
    elif todo.exists():
        targets = [l for l in todo.read_text().splitlines() if l.strip()]
    else:
        print("[extract] no todo list and no --all; nothing to do")
        return 0

    if not targets:
        print("[extract] nothing new")
        return 0

    todo_summarize = ROOT / "scripts" / ".todo_summarize.txt"
    succeeded = []
    for rel in targets:
        if process(rel):
            succeeded.append(str(Path(rel).with_suffix(".md")))
    todo_summarize.write_text("\n".join(succeeded))
    print(f"[extract] done, {len(succeeded)} files ready for summary")
    return 0


if __name__ == "__main__":
    sys.exit(main())
